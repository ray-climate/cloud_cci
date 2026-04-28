"""Build the cloud_cci progress slide deck (PPTX).

Generates two helper figures into ``slides/_assets/`` (framework block diagram
and a one-month ATLID-SLSTR colocation coverage map) and writes
``slides/cloud_cci_progress_2026-04.pptx``.
"""

from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Pt

ROOT = Path("/gws/pw/j07/nceo_aerosolfire/rsong/project/cloud_cci")
SLIDES = ROOT / "slides"
ASSETS = SLIDES / "_assets"
ASSETS.mkdir(parents=True, exist_ok=True)

# 16:9 deck — width 13.333", height 7.5"
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)

NAVY = RGBColor(0x0B, 0x3D, 0x91)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
ACCENT = RGBColor(0xC8, 0x10, 0x2E)


# ---------------------------------------------------------------------------
# Helper figures
# ---------------------------------------------------------------------------

def make_framework_diagram(out: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 5)
    ax.axis("off")

    def box(x, y, w, h, text, fc, ec="#222", fontsize=11, weight="bold"):
        ax.add_patch(mpatches.FancyBboxPatch(
            (x, y), w, h, boxstyle="round,pad=0.05,rounding_size=0.15",
            facecolor=fc, edgecolor=ec, linewidth=1.2,
        ))
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
                fontsize=fontsize, fontweight=weight, color="#111")

    def arrow(x1, y1, x2, y2):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", lw=1.6, color="#333"))

    # Left column: ORAC inputs
    box(0.2, 3.6, 2.6, 1.0, "ORAC L2\n(SEVIRI / SLSTR)", "#F4D35E")
    box(0.2, 1.6, 2.6, 1.0, "EarthCARE L2\n(ATLID, CPR, AM/AC/ACM)", "#9BC1BC")

    # Centre: data access + collocation
    box(3.4, 3.6, 2.8, 1.0, "earthcare/\nMAAP STAC client", "#E0E0E0")
    box(3.4, 1.6, 2.8, 1.0, "orbit/colocation\nATLID ↔ SLSTR", "#E0E0E0")

    # Validation core
    box(6.8, 2.6, 2.4, 1.4, "validation/\nreaders → collocate\n→ statistics", "#F2A65A")

    # Outputs
    box(9.6, 3.6, 1.3, 1.0, "Track\nfigures", "#CFE0E8")
    box(9.6, 1.6, 1.3, 1.0, "Compare\nstats / plots", "#CFE0E8")

    # Arrows
    arrow(2.8, 4.1, 3.4, 4.1)
    arrow(2.8, 2.1, 3.4, 2.1)
    arrow(6.2, 4.1, 6.8, 3.6)
    arrow(6.2, 2.1, 6.8, 2.8)
    arrow(9.2, 3.6, 9.6, 4.1)
    arrow(9.2, 2.6, 9.6, 2.1)

    fig.tight_layout()
    fig.savefig(out, dpi=180, bbox_inches="tight")
    plt.close(fig)


def make_colocation_map(out: Path) -> None:
    csv_dir = ROOT / "orbit" / "colocation"
    months = sorted(csv_dir.glob("atlid_slstr3a_matches_*.csv"))
    counts: list[tuple[str, int]] = []
    for f in months:
        with f.open() as fh:
            n = sum(1 for _ in fh) - 1
        counts.append((f.stem.split("_")[-1], n))

    sample_csv = csv_dir / "atlid_slstr3a_matches_2025-01.csv"
    df = pd.read_csv(sample_csv)
    if len(df) > 6000:
        df = df.sample(6000, random_state=0)

    fig, (ax_map, ax_bar) = plt.subplots(
        1, 2, figsize=(13, 4.6),
        gridspec_kw={"width_ratios": [2.2, 1.0]},
    )

    ax_map.scatter(df["longitude"], df["latitude"], s=1.5, c="#0B3D91",
                   alpha=0.35, linewidths=0)
    ax_map.set_xlim(-180, 180)
    ax_map.set_ylim(-90, 90)
    ax_map.set_xticks([-180, -120, -60, 0, 60, 120, 180])
    ax_map.set_yticks([-90, -60, -30, 0, 30, 60, 90])
    ax_map.set_xlabel("Longitude")
    ax_map.set_ylabel("Latitude")
    ax_map.set_title("ATLID ↔ SLSTR-3A matches — Jan 2025 (sample)")
    ax_map.grid(True, linestyle=":", alpha=0.4)

    labels = [m for m, _ in counts]
    values = [n / 1000 for _, n in counts]
    ax_bar.barh(range(len(labels)), values, color="#0B3D91")
    ax_bar.set_yticks(range(len(labels)))
    ax_bar.set_yticklabels(labels, fontsize=8)
    ax_bar.invert_yaxis()
    ax_bar.set_xlabel("Matches per month (× 1000)")
    ax_bar.set_title("Monthly match counts (S3A)")
    ax_bar.grid(True, axis="x", linestyle=":", alpha=0.4)

    fig.tight_layout()
    fig.savefig(out, dpi=170, bbox_inches="tight")
    plt.close(fig)


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(1.6))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    tf = bar.text_frame
    tf.margin_left = Cm(0.7)
    tf.margin_top = Cm(0.25)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)


def add_footer(slide, page: str) -> None:
    tb = slide.shapes.add_textbox(Cm(0.5), SLIDE_H - Cm(0.9),
                                  SLIDE_W - Cm(1.0), Cm(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Cloud_cci progress — ORAC × EarthCARE validation   •   Oxford / NCEO   •   {page}"
    r.font.size = Pt(9)
    r.font.color.rgb = GREY


def add_textbox(slide, x, y, w, h, lines, *, size=14, bullet=True, bold=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = (f"•  {line}" if bullet else line)
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        if bold:
            r.font.bold = True
        p.space_after = Pt(4)
    return tb


def add_image_caption(slide, x, y, w, text):
    tb = slide.shapes.add_textbox(x, y, w, Cm(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = GREY


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build() -> Path:
    framework_png = ASSETS / "framework.png"
    coloc_png = ASSETS / "colocation_coverage.png"
    make_framework_diagram(framework_png)
    make_colocation_map(coloc_png)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    fig_dir = ROOT / "figures"

    # ------------------------------------------------------------------
    # 1. Title
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    tb = s.shapes.add_textbox(Cm(2.0), Cm(5.5), SLIDE_W - Cm(4.0), Cm(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Validating ORAC cloud retrievals with EarthCARE"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = "Framework progress — cloud_cci meeting"
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

    p3 = tf.add_paragraph()
    p3.space_before = Pt(24)
    r3 = p3.add_run()
    r3.text = "Rui Song   •   University of Oxford / NCEO"
    r3.font.size = Pt(16)
    r3.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

    p4 = tf.add_paragraph()
    r4 = p4.add_run()
    r4.text = "April 2026   •   ESA cloud_cci"
    r4.font.size = Pt(14)
    r4.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # ------------------------------------------------------------------
    # 2. Objectives & scope
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Objectives & scope",
                  "Independent EarthCARE references for ORAC L2 cloud properties")
    add_textbox(s, Cm(1.2), Cm(2.4), Cm(31), Cm(13), [
        "Goal: build a reproducible validation framework for ORAC cloud retrievals against EarthCARE L2 products.",
        "Variables in scope: cloud optical thickness (cot), effective radius (cer), top height (cth), water path (cwp), thermodynamic phase.",
        "Reference instruments: active sensors (ATLID lidar, CPR radar) and active+passive synergy (AM, AC, ACM) products only.",
        "Out of scope: MSI-only passive products (M-COP, M-CM, M-AOT) — too close to ORAC's own measurement domain to be independent.",
        "Deliverables for this milestone: data access module, collocation grid, statistics pipeline, first inter-comparison results.",
    ], size=16)
    add_footer(s, "2 / 12")

    # ------------------------------------------------------------------
    # 3. Framework overview
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Validation framework — overview")
    s.shapes.add_picture(str(framework_png), Cm(1.2), Cm(3.0),
                         width=Cm(31.5))
    add_textbox(s, Cm(1.2), Cm(15.0), Cm(31.5), Cm(2.5), [
        "Three layers: (i) ORAC + EarthCARE inputs, (ii) data access & space-time collocation, (iii) statistics and figures.",
        "All layers are now in place; the EarthCARE → ORAC statistics layer is the next focus (slide 11).",
    ], size=13)
    add_footer(s, "3 / 12")

    # ------------------------------------------------------------------
    # 4. Reference products
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "EarthCARE reference products",
                  "Active and synergy L2; primaries highlighted")

    rows = [
        ("Code", "MAAP type", "Instruments", "Validates ORAC"),
        ("A-CTH", "ATL_CTH_2A", "ATLID", "cth"),
        ("A-TC", "ATL_TC__2A", "ATLID", "phase, cldtype, mask"),
        ("A-EBD", "ATL_EBD_2A", "ATLID", "cot (∫ extinction)"),
        ("A-ICE", "ATL_ICE_2A", "ATLID", "ice cer, cwp"),
        ("C-CLD", "CPR_CLD_2A", "CPR", "cwp, cer"),
        ("C-TC", "CPR_TC__2A", "CPR", "cth, phase"),
        ("AM-CTH ★", "AM__CTH_2B", "ATLID + MSI", "cth across MSI swath"),
        ("AC-TC", "AC__TC__2B", "ATLID + CPR", "phase, cldtype, mask"),
        ("ACM-CAP ★", "ACM_CAP_2B", "ATLID + CPR + MSI", "cot, cer, cwp, phase"),
    ]
    n_rows, n_cols = len(rows), len(rows[0])
    table_shape = s.shapes.add_table(n_rows, n_cols,
                                     Cm(1.2), Cm(2.4),
                                     Cm(31.5), Cm(11.5))
    table = table_shape.table
    widths_cm = [4.0, 6.0, 9.0, 12.5]
    for i, w in enumerate(widths_cm):
        table.columns[i].width = Cm(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = val
            if ri == 0:
                r.font.bold = True
                r.font.size = Pt(14)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                r.font.size = Pt(12)
                if "★" in row[0]:
                    r.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xCC)

    add_textbox(s, Cm(1.2), Cm(14.5), Cm(31.5), Cm(3.0), [
        "Primary references: AM-CTH (cloud-top across the MSI swath) and ACM-CAP (flagship variational retrieval).",
        "ATLID/CPR L2a serve as nadir-only fall-back references and for cross-checking the synergy products.",
    ], size=13)
    add_footer(s, "4 / 12")

    # ------------------------------------------------------------------
    # 5. earthcare/ data access module
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "earthcare/ — ESA MAAP data-access module")
    add_textbox(s, Cm(1.2), Cm(2.4), Cm(31.5), Cm(13), [
        "Python client for the ESA MAAP STAC catalogue (EarthCAREL2Validated_MAAP).",
        "Catalogue search is anonymous; downloads use a 90-day offline token (env var or ~/.maap/offline_token), exchanged automatically for short-lived access tokens.",
        "Per-product spatial / temporal search; output layout: earthcare_data/<PRODUCT_TYPE>/YYYY/MM/DD/<file>.h5.",
        "Both Python API and CLI: `python -m earthcare {list-products | describe | search | download}`.",
        "Currently mirrored locally: ATL_CTH_2A, ATL_EBD_2A, AM__CTH_2B, ACM_CAP_2B.",
    ], size=15)

    # code box
    code_tb = s.shapes.add_textbox(Cm(1.2), Cm(13.0), Cm(31.5), Cm(4.5))
    code_tb.fill.solid()
    code_tb.fill.fore_color.rgb = LIGHT
    code_tf = code_tb.text_frame
    code_tf.word_wrap = True
    snippet = [
        "from earthcare import MaapCatalog",
        "cat = MaapCatalog()",
        "items = cat.search('AM-CTH', start='2025-02-01T00:00:00Z',",
        "                   end='2025-02-02T00:00:00Z', bbox=(-30,-60,60,60))",
        "for it in items[:3]: cat.download_item(it, 'earthcare_data')",
    ]
    for i, line in enumerate(snippet):
        p = code_tf.paragraphs[0] if i == 0 else code_tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = "Courier New"
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    add_footer(s, "5 / 12")

    # ------------------------------------------------------------------
    # 6. Example EarthCARE products
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "EarthCARE products — example frames")

    panel_w = Cm(10.5)
    y = Cm(2.6)
    s.shapes.add_picture(str(fig_dir / "earthcare/a_cth_nadir.png"),
                         Cm(0.6), y, width=panel_w)
    s.shapes.add_picture(str(fig_dir / "earthcare/am_cth_frame.png"),
                         Cm(11.6), y, width=panel_w)
    s.shapes.add_picture(str(fig_dir / "earthcare/acm_cap_curtain.png"),
                         Cm(22.6), y, width=panel_w)

    cap_y = Cm(11.5)
    add_image_caption(s, Cm(0.6), cap_y, panel_w,
                      "A-CTH — ATLID nadir cloud-top height")
    add_image_caption(s, Cm(11.6), cap_y, panel_w,
                      "AM-CTH — ATLID+MSI synergy CTH (frame)")
    add_image_caption(s, Cm(22.6), cap_y, panel_w,
                      "ACM-CAP — active+passive curtain")

    add_textbox(s, Cm(1.2), Cm(13.0), Cm(31.5), Cm(4.5), [
        "Generated by scripts/make_earthcare_test_figures.py from locally mirrored frames.",
        "Confirms readers / projection / variable extraction for both nadir (ATLID) and synergy (AM, ACM) products.",
        "These three are the templates for ORAC ↔ EarthCARE comparison plots in the next milestone.",
    ], size=13)
    add_footer(s, "6 / 12")

    # ------------------------------------------------------------------
    # 7. ATLID-SLSTR collocation grid
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "ATLID ↔ SLSTR collocation grid",
                  "12 months × S3A & S3B — computed on JASMIN via SLURM")

    s.shapes.add_picture(str(coloc_png), Cm(0.8), Cm(2.4), width=Cm(32))
    add_textbox(s, Cm(1.2), Cm(13.0), Cm(31.5), Cm(4.5), [
        "ATLID nadir track (sampled central points) matched to per-orbit SLSTR-3A / 3B swath polygons.",
        "Match criterion: lat/lon point-in-polygon plus a configurable time window; closest-in-time match retained.",
        "Output: orbit/colocation/atlid_slstr3{a,b}_matches_YYYY-MM.csv  (one file per month per platform).",
        "Aug-2024 → Jul-2025 covered for both S3A and S3B; ~10–20k matches per month per platform.",
    ], size=13)
    add_footer(s, "7 / 12")

    # ------------------------------------------------------------------
    # 8. validation/ pipeline
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "validation/ — pipeline modules")

    add_textbox(s, Cm(1.2), Cm(2.4), Cm(15.5), Cm(13), [
        "readers.py — ORAC L2 + EarthCARE H5 readers; common record schema.",
        "collocate.py — swath ↔ swath / nadir matching (uses the orbit/colocation tables).",
        "reference.py — wraps the EarthCARE reference variables per product.",
        "statistics.py — bias, RMSE, MAE, correlation, slope/intercept, by surface and QC mode.",
        "compare_figures.py — R10 vs R11 inter-comparison plots and CSV.",
        "track_figures.py — per-track case-study panels with coastlines / borders / ice mask.",
        "water_cloud_figures.py / cth_figures.py — variable-specific spatial maps.",
    ], size=13)

    # right column: CLI usage
    code_tb = s.shapes.add_textbox(Cm(17.2), Cm(2.6), Cm(15.5), Cm(13))
    code_tb.fill.solid()
    code_tb.fill.fore_color.rgb = LIGHT
    code_tf = code_tb.text_frame
    code_tf.word_wrap = True
    snippet = [
        "# CLI entry point",
        "python -m validation \\",
        "    compare --variable cot_water \\",
        "    --month 2026-02 \\",
        "    --retrievals R10 R11 \\",
        "    --qc strict",
        "",
        "python -m validation \\",
        "    track --frame 09865H \\",
        "    --label marine-stratocumulus",
    ]
    for i, line in enumerate(snippet):
        p = code_tf.paragraphs[0] if i == 0 else code_tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = "Courier New"
        r.font.size = Pt(12)
        r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    add_footer(s, "8 / 12")

    # ------------------------------------------------------------------
    # 9. R10 vs R11 inter-comparison
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Inter-comparison — R10 vs R11 (cot_water, Feb 2026)",
                  "Pixel-level, strict QC")

    panel_w = Cm(10.0)
    y_top = Cm(2.4)
    s.shapes.add_picture(
        str(fig_dir / "cot_water_2026-02_compare/compare_R10_R11_scatter_pixel.png"),
        Cm(0.6), y_top, width=panel_w)
    s.shapes.add_picture(
        str(fig_dir / "cot_water_2026-02_compare/compare_R10_R11_bias_pixel.png"),
        Cm(11.4), y_top, width=panel_w)
    s.shapes.add_picture(
        str(fig_dir / "cot_water_2026-02_compare/compare_R10_R11_rmse_pixel.png"),
        Cm(22.2), y_top, width=panel_w)

    cap_y = Cm(8.4)
    add_image_caption(s, Cm(0.6), cap_y, panel_w, "Scatter (R10 / R11 vs reference)")
    add_image_caption(s, Cm(11.4), cap_y, panel_w, "Bias map")
    add_image_caption(s, Cm(22.2), cap_y, panel_w, "RMSE map")

    add_textbox(s, Cm(1.2), Cm(10.5), Cm(31.5), Cm(7.0), [
        "All-pixel summary (qc_strict): R10  bias = +5.60, RMSE = 40.4   |   R11  bias = +6.93, RMSE = 43.7   (n ≈ 4.5×10⁵).",
        "R11 currently shows slightly higher bias and RMSE than R10 in cot_water at pixel level — worth flagging for the cot_water surface stratification (ocean / land / coast) shown in the by-surface panels.",
        "These plots use the existing ORAC reference; the EarthCARE (ACM-CAP) reference will replace it as the canonical truth in the next milestone.",
    ], size=13)
    add_footer(s, "9 / 12")

    # ------------------------------------------------------------------
    # 10. Case-study tracks
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Case-study tracks — R10 vs R11",
                  "10 frames sampled across regimes; 3 shown")

    panel_w = Cm(10.5)
    y = Cm(2.4)
    s.shapes.add_picture(
        str(fig_dir / "cot_water_2026-02_track_studies/track_09865H_marine-stratocumulus_R10_vs_R11.png"),
        Cm(0.6), y, width=panel_w)
    s.shapes.add_picture(
        str(fig_dir / "cot_water_2026-02_track_studies/track_09542B_north-atlantic_R10_vs_R11.png"),
        Cm(11.6), y, width=panel_w)
    s.shapes.add_picture(
        str(fig_dir / "cot_water_2026-02_track_studies/track_09574A_tropical-atlantic_R10_vs_R11.png"),
        Cm(22.6), y, width=panel_w)

    cap_y = Cm(13.4)
    add_image_caption(s, Cm(0.6), cap_y, panel_w, "Marine stratocumulus (09865H)")
    add_image_caption(s, Cm(11.6), cap_y, panel_w, "North Atlantic (09542B)")
    add_image_caption(s, Cm(22.6), cap_y, panel_w, "Tropical Atlantic (09574A)")

    add_textbox(s, Cm(1.2), Cm(14.5), Cm(31.5), Cm(3.0), [
        "Track panels include coastlines, country borders, and an ice mask (liquid_classification == 3).",
        "Full set of 10 cases covers Arctic, North Africa, North/Tropical Atlantic, Western Europe, Indian Ocean, S. America, S. Indian Ocean, marine Sc, tropical mixed.",
    ], size=13)
    add_footer(s, "10 / 12")

    # ------------------------------------------------------------------
    # 11. Status & next steps
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Status & next steps")

    # Done column
    head = s.shapes.add_textbox(Cm(1.2), Cm(2.2), Cm(15.5), Cm(1.0))
    p = head.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Completed"
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY
    add_textbox(s, Cm(1.2), Cm(3.4), Cm(15.5), Cm(13), [
        "earthcare/ MAAP STAC client (search + download).",
        "Local mirror: A-CTH, A-EBD, AM-CTH, ACM-CAP test frames.",
        "ATLID ↔ SLSTR collocation grid (12 months × S3A & S3B).",
        "validation/ pipeline (readers, collocate, statistics, figures).",
        "R10 vs R11 inter-comparison — cot_water, cer_water, cth (Feb 2026).",
        "10 case-study tracks across regimes.",
    ], size=13)

    head2 = s.shapes.add_textbox(Cm(17.2), Cm(2.2), Cm(15.5), Cm(1.0))
    p = head2.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Next"
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = ACCENT
    add_textbox(s, Cm(17.2), Cm(3.4), Cm(15.5), Cm(13), [
        "Connect ACM-CAP (cot/cer/cwp/phase) and AM-CTH (cth) as the canonical references inside validation/statistics.py.",
        "Bulk-download EarthCARE frames over the SEVIRI / SLSTR disk for the colocation period (Aug-2024 → Jul-2025).",
        "Extend collocation from ATLID-only to MSI footprint matching for AM-CTH / ACM-CAP comparisons.",
        "Add ice-cloud cases (A-ICE, C-CLD) to the case-study set.",
        "Automate monthly statistics and version regression checks.",
    ], size=13)

    add_footer(s, "11 / 12")

    # ------------------------------------------------------------------
    # 12. Summary / questions
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    tb = s.shapes.add_textbox(Cm(2.0), Cm(4.0), SLIDE_W - Cm(4.0), Cm(11.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Summary"
    r.font.size = Pt(32); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    bullets = [
        "Validation framework end-to-end in place: data access, collocation, statistics, figures.",
        "First inter-comparison results (R10 vs R11) reproducible from the pipeline.",
        "Next milestone: swap the reference layer from ORAC-internal to EarthCARE (AM-CTH / ACM-CAP).",
    ]
    for line in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run(); r.text = "•  " + line
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

    p = tf.add_paragraph()
    p.space_before = Pt(40)
    r = p.add_run(); r.text = "Questions / discussion"
    r.font.size = Pt(28); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xD9, 0x66)

    out = SLIDES / "cloud_cci_progress_2026-04.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    out = build()
    print(f"wrote {out}")
