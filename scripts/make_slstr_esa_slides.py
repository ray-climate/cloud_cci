"""Build the ESA-workshop deck (.pptx) for the SLSTR x EarthCARE validation.

Clean, science-focused academic style: Arial throughout, white background, a thin
accent rule under each title (no filled colour blocks), figure-forward slides with
tight supporting bullets. Fully editable in PowerPoint / Keynote / LibreOffice.

Run: .venv_slides/bin/python scripts/make_slstr_esa_slides.py
Out: docs/presentations/SLSTR_EarthCARE_validation_ESA.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- style ----
FONT = "Arial"
INK = RGBColor(0x1A, 0x1A, 0x1A)   # near-black body text
ACCENT = RGBColor(0x14, 0x4E, 0x74)  # restrained deep steel-blue
MUTED = RGBColor(0x5A, 0x5A, 0x5A)   # subtitles / captions
FAINT = RGBColor(0x8A, 0x8A, 0x8A)   # footer
RULEFILL = RGBColor(0x14, 0x4E, 0x74)
ROWALT = RGBColor(0xEF, 0xF3, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ------------------------------------------------ editable title metadata ----
PRESENTER = "[Presenter name]"
AFFIL = "[Affiliation]"
WORKSHOP = "[ESA workshop name]"
DATE = "[Month 2026]"

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "presentations" / "SLSTR_EarthCARE_validation_ESA.pptx"
F = ROOT / "figures"

SW, SH = Inches(13.333), Inches(7.5)
RUNNING = "ORAC SLSTR × EarthCARE validation"


# ------------------------------------------------------------- helpers ----
def _font(run, size, *, bold=False, italic=False, color=INK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _fit(path, max_w_in, max_h_in):
    with Image.open(path) as im:
        w, h = im.size
    ar = w / h
    if max_w_in / ar <= max_h_in:
        return max_w_in, max_w_in / ar
    return max_h_in * ar, max_h_in


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def footer(slide, page):
    lb = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(9.0), Inches(0.3))
    r = lb.text_frame.paragraphs[0].add_run()
    r.text = f"{RUNNING}  ·  {WORKSHOP}"
    _font(r, 9, color=FAINT)
    rb = slide.shapes.add_textbox(Inches(12.2), Inches(7.12), Inches(0.8), Inches(0.3))
    p = rb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(page)
    _font(r, 9, color=FAINT)


def title_bar(slide, text, sub=None, page=None):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.30), Inches(12.2), Inches(0.62))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    _font(r, 24, bold=True, color=INK)
    # thin accent rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.05),
                                  Inches(12.2), Pt(2.2))
    rule.fill.solid(); rule.fill.fore_color.rgb = RULEFILL
    rule.line.fill.background(); rule.shadow.inherit = False
    if sub:
        sbx = slide.shapes.add_textbox(Inches(0.55), Inches(1.12), Inches(12.2), Inches(0.35))
        r = sbx.text_frame.paragraphs[0].add_run(); r.text = sub
        _font(r, 13, italic=True, color=MUTED)
    if page is not None:
        footer(slide, page)


def bullets(slide, points, *, left=0.65, top=5.32, width=12.05, height=1.95, size=14):
    """points: list of (text, level, lead_bool)."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, (text, level, lead) in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(5)
        p.line_spacing = 1.06
        marker = "– " if level else "▪ "
        mrun = p.add_run(); mrun.text = marker
        _font(mrun, size, bold=lead, color=ACCENT if not level else MUTED)
        trun = p.add_run(); trun.text = text
        _font(trun, size, bold=lead, color=INK)


def add_image(slide, path, *, top, max_w, max_h, left=None):
    w, h = _fit(path, max_w, max_h)
    if left is None:
        left = (13.333 - w) / 2.0
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(w), Inches(h))


def fig_slide(prs, title, sub, fig, points, page, *, bullets_top=5.32):
    s = blank(prs)
    title_bar(s, title, sub, page)
    add_image(s, fig, top=1.5, max_w=11.8, max_h=3.55)
    bullets(s, points, top=bullets_top)
    return s


def two_fig_slide(prs, title, sub, fig_l, fig_r, points, page):
    s = blank(prs)
    title_bar(s, title, sub, page)
    add_image(s, fig_l, top=1.55, max_w=6.0, max_h=3.45, left=0.55)
    add_image(s, fig_r, top=1.55, max_w=6.0, max_h=3.45, left=6.78)
    bullets(s, points, top=5.32)
    return s


# --------------------------------------------------------------- deck ----
def main():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # -- 1. Title -------------------------------------------------------
    s = blank(prs)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "Validating ORAC SLSTR cloud retrievals against EarthCARE"
    _font(r, 32, bold=True, color=INK)
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "Sentinel-3A × EarthCARE — December 2025"
    _font(r, 18, color=ACCENT)
    # accent rule
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.92), Inches(3.95),
                              Inches(4.2), Pt(2.5))
    rule.fill.solid(); rule.fill.fore_color.rgb = RULEFILL
    rule.line.fill.background(); rule.shadow.inherit = False
    # provenance
    cap = s.shapes.add_textbox(Inches(0.92), Inches(4.15), Inches(11.5), Inches(1.0))
    ct = cap.text_frame; ct.word_wrap = True
    r = ct.paragraphs[0].add_run()
    r.text = ("Continuous: CTH · water/ice COT · CER · CWP     "
              "Categorical: cloud mask · phase")
    _font(r, 14, color=MUTED)
    r2 = ct.add_paragraph().add_run()
    r2.text = "References: EarthCARE ATLID (A-CTH, A-EBD, A-TC) + ACM-CAP synergy"
    _font(r2, 14, color=MUTED)
    # author block
    ab = s.shapes.add_textbox(Inches(0.92), Inches(5.75), Inches(11.5), Inches(1.1))
    at = ab.text_frame; at.word_wrap = True
    r = at.paragraphs[0].add_run(); r.text = PRESENTER
    _font(r, 15, bold=True, color=INK)
    r = at.add_paragraph().add_run(); r.text = AFFIL
    _font(r, 13, color=MUTED)
    r = at.add_paragraph().add_run(); r.text = f"{WORKSHOP}  ·  {DATE}"
    _font(r, 13, color=MUTED)

    # -- 2. Motivation & objective -------------------------------------
    s = blank(prs)
    title_bar(s, "Motivation & objective",
              "Why validate ORAC SLSTR cloud products against an active reference", 2)
    bullets(s, [
        ("ORAC retrieves cloud properties (CTH, COT, CER, CWP, phase, mask) from the SLSTR dual-view imager on Sentinel-3A.", 0, True),
        ("Passive cloud retrieval is hardest over the bright, cold cryosphere — this ORAC build (v5.1_new_snowice) specifically revises snow / sea-ice surface handling.", 0, False),
        ("EarthCARE supplies an independent active reference (ATLID lidar + CPR radar + MSI synergy) whose physics ORAC does not share — so disagreement reveals real retrieval error.", 0, False),
        ("Objective: quantify bias, scatter and skill for every EarthCARE-validatable ORAC SLSTR variable over December 2025 — the one full processed month of this stream.", 0, True),
    ], top=1.7, height=4.8, size=17)

    # -- 3. Reference strategy -----------------------------------------
    s = blank(prs)
    title_bar(s, "Reference strategy — the most independent measurement wins",
              "ATLID active lidar first; MSI passive products never (Holz 2008; Cloud_cci PVIR v6)", 3)
    bullets(s, [
        ("Active lidar / radar shares no physics with a passive imager — a disagreement is a genuine ORAC error, not a shared passive bias.", 0, True),
        ("MSI-derived products are excluded on principle: agreement there could be two passive sensors making the same mistake.", 0, False),
    ], top=1.55, height=1.3, size=15)
    rows = [
        ("ORAC SLSTR variable", "EarthCARE reference", "Independence"),
        ("Cloud-top height (cth)", "A-CTH  (ATLID)", "pure lidar"),
        ("Water COT / CER / CWP", "ACM-CAP  (ATLID+CPR+MSI)", "radar+lidar synergy"),
        ("Ice COT", "A-EBD  (ATLID ∫α dz)", "pure lidar"),
        ("Cloud phase & mask", "A-TC  (ATLID target class.)", "pure lidar"),
    ]
    tbl = s.shapes.add_table(len(rows), 3, Inches(0.65), Inches(3.05),
                             Inches(12.0), Inches(3.1)).table
    tbl.columns[0].width = Inches(4.1); tbl.columns[1].width = Inches(4.7)
    tbl.columns[2].width = Inches(3.2)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            run = para.add_run(); run.text = val
            _font(run, 13 if ri else 13, bold=(ri == 0 or ci == 0),
                  color=WHITE if ri == 0 else INK)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROWALT if ri % 2 else WHITE

    # -- 4. Collocation geometry ---------------------------------------
    fig_slide(prs,
        "Collocation — a polar comparison by orbital mechanics",
        "Two sun-synchronous orbiters coincide only near their track crossings",
        F / "slstr_collocation" / "collocation_map_polar.png",
        [("EarthCARE and Sentinel-3A are both sun-synchronous at different local overpass times → simultaneous views only where the orbit planes converge, near the poles.", 0, True),
         ("1.71 M matches, every one at |lat| 70.6–83.0° (853k Arctic, 860k Antarctic); tropics and mid-latitudes return N = 0 — orbital geometry, not a filter.", 0, False),
         ("Complementary to the geostationary SEVIRI validation (±60°): SLSTR covers the poles against the same ATLID truth.", 0, False)],
        4)

    # -- 5. Anatomy of a match -----------------------------------------
    two_fig_slide(prs,
        "Anatomy of a match — sub-pixel, well inside the thresholds",
        "One ATLID profile paired to the SLSTR pixel it falls in (frame 08642G, Antarctica)",
        F / "slstr_collocation" / "crossing_case_study.png",
        F / "slstr_collocation" / "match_quality.png",
        [("Temporal gate |Δt| ≤ 60 min (median 26 min); spatial nearest-pixel ≤ 3 km on-swath gate (median 0.43 km, max ~1.1 km).", 0, True),
         ("SLSTR ≈ EarthCARE ≈ 1 km footprint (measured) → matches are sub-pixel for both; negligible footprint mismatch, unlike the 3–7 km SEVIRI pixels.", 0, False),
         ("The 3 km gate is a swath-membership test, not an averaging radius — it sits far above the 0.43 km match median.", 0, False)],
        5)

    # -- 6. Thresholds non-binding -------------------------------------
    two_fig_slide(prs,
        "Both thresholds are non-binding — robustness",
        "CTH statistics stay flat across the temporal and spatial windows",
        F / "slstr_dt_sweep" / "slstr_dt_sweep.png",
        F / "slstr_sensitivity" / "slstr_cth_sensitivity.png",
        [("Across Δt {15, 30, 45, 60} min × distance {1, 2, 3} km: CTH bias −0.55 to −0.57 km, RMSE 2.08–2.11 km, R 0.58–0.60 — essentially unchanged.", 0, True),
         ("Polar clouds evolve slowly and matches are already < 1 km, so neither threshold tunes the result.", 0, False),
         ("Neither can be relaxed toward lower latitudes without breaking the comparison (hours in Δt, or hundreds of km in distance).", 0, False)],
        6)

    # -- 7. CTH ---------------------------------------------------------
    two_fig_slide(prs,
        "Cloud-top height — the strong result",
        "ORAC cth_corrected vs A-CTH; thermal retrieval works day and night",
        F / "slstr_cth_2025-12" / "cth_scatter.png",
        F / "slstr_cth_2025-12" / "cth_by_cloud_type.png",
        [("Median bias −0.12 km (mean −0.57, pulled by the multi-layer tail), RMSE 2.08 km, R 0.58 (N = 162k) — the typical cloud top is essentially unbiased and correlated.", 0, True),
         ("Thick single-layer cloud is near-perfect (median −0.02 km, R 0.75); the error concentrates in high / multi-layer cloud (≈ −4 km) — the classic passive multi-layer ambiguity.", 0, False),
         ("Hemispheric contrast: Arctic −0.90 km vs Antarctic −0.28 km (December = Arctic night / Antarctic day).", 0, False)],
        7)

    # -- 8. Water COT saturation ---------------------------------------
    fig_slide(prs,
        "Water-cloud optical depth — the passive retrieval saturates",
        "ORAC vs ACM-CAP liquid τ (radar-aided synergy), phase-agree liquid",
        F / "slstr_cot_water_2025-12" / "cot_water_saturation.png",
        [("Report the median: −4.8. The +3.1 mean is a heavy-tail skew artefact — it even flips sign.", 0, True),
         ("ORAC passive liquid τ is pinned ~7–8 across the entire ACM-CAP range (0.6 → 34): with no dynamic range it cannot correlate (r_log 0.11).", 0, False),
         ("It is ORAC saturating, not the radar reference reading high (adding CPR changes the bias by ~1 τ).", 0, False)],
        8)

    # -- 9. Surface split ----------------------------------------------
    fig_slide(prs,
        "The deficit is cryospheric — the surface, not the phase",
        "Solar-retrieval skill split by ORAC surface type",
        F / "slstr_surface_2025-12" / "surface_type_bias.png",
        [("Water COT correlates (r = +0.28) and over-reads only over open water; over sea-ice & ice-sheet (95% of the scene) r → 0 and the bias is −4 to −5.", 0, True),
         ("Same clouds, only the background changed → the saturation is the bright-surface radiative regime, not an algorithm bug.", 0, False),
         ("CER stays surface-robust (median +0.2 / +0.3 µm over sea-ice / ice-sheet): trustworthy droplet size even where τ is unconstrained.", 0, False)],
        9)

    # -- 10. Ice COT & CER ---------------------------------------------
    two_fig_slide(prs,
        "Ice optical depth & effective radius — low bias, weak skill",
        "Ice COT vs A-EBD column τ  ·  CER vs ACM-CAP liquid rₑ",
        F / "slstr_cot_ice_2025-12" / "cot_scatter.png",
        F / "slstr_cer_water_2025-12" / "cer_water_scatter.png",
        [("Ice COT median +2.0 (mean +7.2 skewed; RMSE 16, r_log 0.17) — a modest typical overestimate; land ≈ ocean once the high-τ tail is removed.", 0, True),
         ("CER median +1.1 µm — nearly unbiased, but R ≈ −0.1: the SWIR radius scatters around the central value rather than tracking true particle size.", 0, False),
         ("Both are limited by scatter over bright, high-sun-zenith surfaces — not by a large systematic offset.", 0, False)],
        10)

    # -- 11. CWP --------------------------------------------------------
    fig_slide(prs,
        "Liquid water path — the saturation reaches the water budget",
        "ORAC cwp vs an independent ACM-CAP water-content reference (LWP)",
        F / "slstr_cwp_2025-12" / "cwp_validation.png",
        [("Median bias −16 g m⁻² (−34%): 30 vs 47 g m⁻², decorrelated (r_log 0.02), worse over ice sheet (−18) than ocean (−13).", 0, True),
         ("Validated against radar+lidar liquid-water-content — a reference that never sees τ — so it independently confirms the τ saturation propagates into the water budget.", 0, False),
         ("Fixed by the same surface-albedo advance, not a CWP-specific change.", 0, False)],
        11)

    # -- 12. Phase & cloud mask ----------------------------------------
    fig_slide(prs,
        "Cloud mask & phase vs A-TC — the two-way contingency",
        "Categorical validation against ATLID Target Classification (N = 614k pixels)",
        F / "slstr_phase_2025-12" / "phase_contingency.png",
        [("Cloud mask: POD 0.69, FAR 0.11, accuracy 0.79 — conservative; the missed 31% is 88% thin cirrus the lidar sees and a passive imager cannot (the irreducible passive limit).", 0, True),
         ("Phase: POD_liquid 89.5%, POD_ice 62.4% → ORAC calls 38% of ice cloud tops 'liquid' — a liquid bias.", 0, True),
         ("This bias is surface-independent → a second, intrinsic limitation, distinct from the surface-driven τ saturation.", 0, False)],
        12)

    # -- 13. All variables at a glance ---------------------------------
    s = blank(prs)
    title_bar(s, "All variables at a glance — December 2025",
              "Median-primary headline per variable (N = pixel-level matches)", 13)
    rows = [
        ("Variable", "Reference", "Headline metric", "Verdict"),
        ("Cloud-top height", "A-CTH", "median −0.12 km, R 0.58", "Strong"),
        ("Water-cloud COT", "ACM-CAP", "median −4.8 (passive τ saturates ~7–8)", "Regime-limited"),
        ("Ice-cloud COT", "A-EBD", "median +2.0 (r_log 0.17)", "Weak corr."),
        ("Effective radius", "ACM-CAP", "median +1.1 µm (R ≈ −0.1)", "Robust bias, no skill"),
        ("Liquid water path", "ACM-CAP LWP", "median −16 g m⁻² (−34%)", "Inherits τ saturation"),
        ("Cloud mask", "A-TC", "POD 0.69 · FAR 0.11", "Conservative"),
        ("Cloud phase", "A-TC", "POD_liq 90% · POD_ice 62%", "Liquid-biased"),
    ]
    nr, nc = len(rows), 4
    tbl = s.shapes.add_table(nr, nc, Inches(0.55), Inches(1.55),
                             Inches(12.25), Inches(4.9)).table
    tbl.columns[0].width = Inches(2.7); tbl.columns[1].width = Inches(1.9)
    tbl.columns[2].width = Inches(5.15); tbl.columns[3].width = Inches(2.5)
    for ri in range(nr):
        for ci in range(nc):
            cell = tbl.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = cell.text_frame.paragraphs[0].add_run(); run.text = rows[ri][ci]
            _font(run, 14 if ri == 0 else 13, bold=(ri == 0 or ci == 0),
                  color=WHITE if ri == 0 else INK)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROWALT if ri % 2 else WHITE

    # -- 14. Summary & next steps --------------------------------------
    s = blank(prs)
    title_bar(s, "Summary — what is validated, and what is next",
              "Every EarthCARE-validatable dimension of the Dec-2025 sample is complete", 14)
    bullets(s, [
        ("VALIDATED (all median-primary, with confidence intervals):", 0, True),
        ("CTH · water-COT · ice-COT · CER · CWP · cloud mask · phase, plus the collocation methodology and surface-type / phase stratification.", 1, False),
        ("TWO INDEPENDENT LIMITATIONS OF THE POLAR-DAYTIME SOLAR RETRIEVAL:", 0, True),
        ("(1) bright-surface optical-depth SATURATION — surface-driven, propagates into CWP, spares open water;", 1, False),
        ("(2) an intrinsic LIQUID PHASE BIAS — POD_ice 62%, surface-independent.", 1, False),
        ("TRUSTWORTHY OVER THE CRYOSPHERE:", 0, True),
        ("thermal CTH (median −0.12 km, R 0.58) and CER droplet size are robust; the solar τ / water-path products are regime-limited.", 1, False),
        ("NEXT STEPS (require new processing, not further analysis):", 0, True),
        ("an Arctic boreal-summer month · Sentinel-3B · other seasons; and ORAC-SLSTR vs ORAC-SEVIRI against the same ATLID truth.", 1, False),
    ], top=1.55, height=5.4, size=15)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("wrote", OUT, f"({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
