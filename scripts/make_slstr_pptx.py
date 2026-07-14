"""Build a PowerPoint deck (.pptx) for the SLSTR x EarthCARE validation progress
meeting. Embeds the actual analysis figures with headline talking points, so it is
fully editable in PowerPoint / Keynote / LibreOffice.

Run: python scripts/make_slstr_pptx.py
Out: docs/presentations/SLSTR_EarthCARE_validation_Dec2025.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1B, 0x4F, 0x72)
ACCENT = RGBColor(0x2E, 0x6D, 0xA4)
GREY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x22, 0x22, 0x22)

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "presentations" / "SLSTR_EarthCARE_validation_Dec2025.pptx"

SW, SH = Inches(13.333), Inches(7.5)


def _fit(path: Path, max_w_in: float, max_h_in: float):
    """Return (width_emu, height_emu) preserving aspect to fit the box."""
    with Image.open(path) as im:
        w, h = im.size
    ar = w / h
    if max_w_in / ar <= max_h_in:
        return Inches(max_w_in), Inches(max_w_in / ar)
    return Inches(max_h_in * ar), Inches(max_h_in)


def title_bar(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(14); r2.font.color.rgb = GREY; r2.font.italic = True
    # accent underline
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(1.18), Inches(12.3), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def add_image(slide, path: Path, top_in=1.4, max_w_in=8.2, max_h_in=5.4, left_in=None):
    w, h = _fit(path, max_w_in, max_h_in)
    left = Inches(left_in) if left_in is not None else Emu(int((SW - w) / 2))
    slide.shapes.add_picture(str(path), left, Inches(top_in), width=w, height=h)
    return w, h


def bullets(slide, items, left_in, top_in, width_in, height_in, size=15):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, lvl, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl)
        r.font.bold = bold
        r.font.color.rgb = NAVY if bold else DARK
        p.space_after = Pt(6)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_notes(slide, text):
    """Attach speaker notes (presenter view) to a slide."""
    if text:
        slide.notes_slide.notes_text_frame.text = text.strip()


def fig_slide(prs, title, sub, fig, points, note=None):
    """Wide figure on top, bullets below."""
    s = blank(prs)
    title_bar(s, title, sub)
    if fig and fig.exists():
        add_image(s, fig, top_in=1.4, max_w_in=12.4, max_h_in=3.9)
    bullets(s, points, left_in=0.7, top_in=5.5, width_in=12.0, height_in=1.9, size=15)
    set_notes(s, note)
    return s


def main():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    F = ROOT / "figures"

    # ---- 1. Title ----
    s = blank(prs)
    bar = s.shapes.add_shape(1, 0, Inches(2.4), SW, Inches(2.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.1))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "Validating ORAC SLSTR cloud retrievals against EarthCARE"
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph(); r2 = p2.add_run()
    r2.text = "Sentinel-3A × EarthCARE — December 2025"
    r2.font.size = Pt(18); r2.font.color.rgb = RGBColor(0xD5, 0xE4, 0xF0)
    cap = s.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.2))
    ct = cap.text_frame; ct.word_wrap = True
    cp = ct.paragraphs[0]; cr = cp.add_run()
    cr.text = ("Continuous validation: CTH · water/ice COT · CER · CWP    |    "
               "Categorical: cloud mask · phase (A-TC)")
    cr.font.size = Pt(14); cr.font.color.rgb = GREY
    set_notes(s,
        "Framing. We've extended our SEVIRI x EarthCARE validation framework to ORAC "
        "retrievals on Sentinel-3A SLSTR, using EarthCARE as the reference, for "
        "December 2025. The story I'll build to is two independent limitations of the "
        "polar-daytime solar retrieval - a bright-surface optical-depth saturation, and "
        "an intrinsic liquid phase bias - while the thermal cloud-top height and the "
        "effective radius come out well. The sample is about 0.6 million matched pixels.")

    # ---- 2. Progress to date ----
    s = blank(prs)
    title_bar(s, "Progress to date",
              "From collocation framework to a complete Antarctic-summer validation")
    bullets(s, [
        ("Collocation pipeline built & characterised — polar simultaneous-nadir-overpass geometry, ±60 min / 3 km, density map + Δt/distance sensitivity.", 0, True),
        ("Continuous validation complete — CTH (both hemispheres), water-COT, ice-COT, CER, CWP, all against EarthCARE and median-primary with confidence intervals.", 0, True),
        ("Root cause established — passive τ saturation over bright surfaces; surface-type stratification; independent radar+lidar water-path check.", 0, False),
        ("Categorical validation closed — NEW: cloud mask + phase vs A-TC; ice-detection skill (POD_ice) measurable for the first time.", 0, True),
        ("Fully written up & reproducible — six variable reports + summary + this deck; all version-controlled.", 0, False),
        ("Ready for review — remaining extensions need new processing (Arctic month, Sentinel-3B, other seasons), not further analysis of this sample.", 0, False),
    ], left_in=0.7, top_in=1.55, width_in=12.0, height_in=5.4, size=17)
    set_notes(s,
        "Quick status before the detail. The collocation pipeline is built and "
        "characterised. Every continuous variable is validated against EarthCARE. We've "
        "done the root-cause analysis behind the biases. And - new since the last meeting "
        "- we've closed the two-way phase validation using A-TC, which lets us report "
        "ice-detection skill for the first time. All of it is written up as reports plus "
        "this deck, and it's reproducible. What remains needs new data processing - a "
        "summer Arctic month, Sentinel-3B, other seasons - not more analysis of what we have.")

    # ---- 3. Scope & method ----
    fig_slide(prs,
        "Scope & collocation — a polar comparison by orbital mechanics",
        "Simultaneous-nadir-overpass geometry of two sun-synchronous orbiters",
        F / "slstr_collocation" / "collocation_map_polar.png",
        [("SLSTR and EarthCARE are both sun-synchronous → their tracks coincide only near the poles (~70–83°). Every match is polar — not a choice.", 0, True),
         ("Match: temporal gate ±60 min + nearest pixel < 3 km (SLSTR ≈ EarthCARE ≈ 1 km footprint).", 0, False),
         ("Solar retrievals need daylight → December = Antarctic polar day → the COT/CER/CWP/phase sample is 100% Southern Hemisphere; thermal CTH also samples the Arctic (night).", 0, False),
         ("References: A-CTH (height) · ACM-CAP (liquid COT/CER, radar+lidar LWP) · A-EBD (ice COT) · A-TC (phase & cloud mask).", 0, False)],
        note=("Why this is a polar study: two sun-synchronous orbiters only coincide near "
              "the poles, so every collocation is polar - that's orbital mechanics, not a "
              "choice. We match within an hour and 3 km, and both footprints are about a "
              "kilometre. The key consequence: the solar retrievals need daylight, and in "
              "December that's Antarctic polar day, so the entire optical-depth, effective-"
              "radius, water-path and phase sample is Southern Hemisphere. Only the thermal "
              "cloud-top height, which works at night, also sees the Arctic. The map shows "
              "where the crossings concentrate - a dense polar band, sparse elsewhere."))

    # ---- 3. CTH ----
    fig_slide(prs,
        "Cloud-top height — the strong result",
        "ORAC vs A-CTH, both hemispheres (thermal retrieval works day & night)",
        F / "slstr_cth_2025-12" / "cth_scatter.png",
        [("Median CTH bias −0.25 km (Antarctic, polar day) — the thermal cloud-top retrieval survives the polar regime.", 0, True),
         ("Arctic (polar night) −0.95 km: ~4× larger, a night-vs-day + winter-sea-ice regime contrast, not a pole effect.", 0, False),
         ("The one variable that samples both hemispheres; the headline −0.57 km 'polar' number blends the two regimes.", 0, False)],
        note=("Cloud-top height is our strongest result. Over the Antarctic in daylight the "
              "median bias is a quarter of a kilometre - the thermal 11-micron retrieval "
              "holds up in the polar regime. The Arctic, in polar night in December, is "
              "about four times worse at minus 0.95 km; that's a night-versus-day and "
              "winter-sea-ice contrast, not something about the pole itself. Worth flagging: "
              "the often-quoted minus 0.57 km 'polar' number is really a blend of these two "
              "very different regimes, so I prefer to quote them split."))

    # ---- 4. Water-COT saturation ----
    fig_slide(prs,
        "Water-cloud optical depth — the passive retrieval saturates",
        "ORAC vs ACM-CAP (radar-aided synergy), phase-agree liquid",
        F / "slstr_cot_water_2025-12" / "cot_water_saturation.png",
        [("Report the MEDIAN: −4.8. The quoted '+3 mean' is a skew artefact of a high-τ tail (it even flips sign).", 0, True),
         ("ORAC's passive liquid τ is pinned ~5–8 across the entire ACM-CAP range (0.6→34): it over-reads thin cloud, saturates on thick, and cannot correlate (r_log ≈ 0.1).", 0, False),
         ("It is ORAC saturating, not the radar reference being high (CPR changes the bias by ~1 τ).", 0, False)],
        note=("Water-cloud optical depth. First a methodological point: report the median, "
              "not the mean. Optical depth is heavy-tailed, and the mean gets dragged "
              "positive by a few thick-cloud pixels - it even flips sign. On the median, "
              "ORAC underestimates by about 5. The cause is saturation: ORAC's passive "
              "liquid optical depth is stuck around 5 to 8 across the whole reference range, "
              "from below 1 up to 34. So it over-reads thin cloud, saturates on thick cloud, "
              "and can't correlate. And it's genuinely ORAC saturating - using the radar in "
              "the reference changes the answer by only about 1, so the reference isn't the "
              "problem."))

    # ---- 5. Surface type ----
    fig_slide(prs,
        "The deficit is cryospheric — the surface, not the phase",
        "Solar-retrieval skill split by ORAC surface type",
        F / "slstr_surface_2025-12" / "surface_type_bias.png",
        [("Water-COT correlates (r = +0.28) and over-reads ONLY over open water; over sea-ice & ice-sheet (95% of the scene) r → 0 and bias is −4 to −5.", 0, True),
         ("Same clouds, only the background changed → the τ saturation is the bright-surface radiative regime, not an algorithm bug.", 0, False),
         ("CER is surface-robust (median +0.2 / +0.3 µm over sea-ice / ice-sheet): trustworthy droplet size even where τ is unconstrained.", 0, False)],
        note=("This is the key slide for the cause. If we split by surface type, water "
              "optical depth only behaves like a working retrieval - positive correlation, "
              "slight over-read - over open water. The instant the surface is sea-ice or "
              "ice-sheet, which is 95% of the scene, the correlation collapses to zero and "
              "it under-reads by 4 to 5. Same clouds, only the background changed - so the "
              "saturation is a bright-surface radiative effect, not a bug. And note the "
              "contrast: effective radius stays accurate over the cryosphere, so droplet "
              "size is trustworthy even where optical depth isn't."))

    # ---- 6. CWP ----
    fig_slide(prs,
        "Liquid water path — the saturation reaches the water budget",
        "ORAC cwp vs an independent ACM-CAP water-content reference (LWP)",
        F / "slstr_cwp_2025-12" / "cwp_validation.png",
        [("ORAC liquid water path runs 34% low on the median (30 vs 47 g m⁻²), decorrelated, worse over ice sheet (−18) than ocean (−13).", 0, True),
         ("Validated against radar+lidar liquid-water-content — a reference that never sees τ — so it confirms the τ saturation propagates into the water budget.", 0, False),
         ("Fixed by the same surface-albedo advance, not a CWP-specific change.", 0, False)],
        note=("Does the optical-depth problem matter for the water budget? Yes. ORAC's "
              "liquid water path runs about a third low on the median, decorrelated, and "
              "worse over the ice sheet. What makes this convincing is the reference: we "
              "integrated EarthCARE's radar-plus-lidar liquid water content, a measurement "
              "that never uses optical depth - so it independently confirms the saturation "
              "propagates into the water path. Same fix required: better handling of the "
              "bright surface albedo."))

    # ---- 7. Phase & cloud mask ----
    fig_slide(prs,
        "Cloud mask & phase vs EarthCARE A-TC — the two-way contingency",
        "Categorical validation (POD_ice, previously unmeasurable)",
        F / "slstr_phase_2025-12" / "phase_contingency.png",
        [("Cloud mask: POD 0.69, FAR 0.11 — conservative; the missed 31% is 88% thin cirrus the lidar sees and a passive imager cannot.", 0, True),
         ("Phase: POD_liquid 89.5%, POD_ice 62.4% → ORAC has a liquid bias, calling 38% of ice cloud tops 'liquid'.", 0, True),
         ("This bias is SURFACE-INDEPENDENT (62–64% over open water, sea-ice, ice-sheet) — an intrinsic limitation, distinct from the surface-driven τ saturation.", 0, False),
         ("Polar liquid is ~100% supercooled; ORAC handles it well (89.5%).", 0, False)],
        note=("The categorical validation, and the part that's new since last time. Against "
              "A-TC, which classifies every lidar bin, we can finally score both phases. "
              "Liquid detection is good at about 90%, but ice detection is only 62% - ORAC "
              "calls 38% of ice cloud tops liquid. That's a liquid bias. Crucially it's "
              "surface-independent - 62 to 64% over open water, sea-ice and ice-sheet alike "
              "- so it's a different, intrinsic limitation, separate from the surface-driven "
              "saturation. On the cloud mask, ORAC is conservative: it catches 69% with a "
              "low false-alarm rate, and what it misses is 88% thin cirrus that the lidar "
              "sees and a passive imager physically cannot. Read POD as: of all the real ice "
              "clouds, what fraction did ORAC catch."))

    # ---- 8. All variables at a glance (table) ----
    s = blank(prs)
    title_bar(s, "All variables at a glance — December 2025",
              "Headline metric per variable (median-primary; N = pixel-level matches)")
    rows = [
        ("Variable", "EarthCARE ref.", "Headline metric", "Verdict"),
        ("Cloud-top height", "A-CTH", "median −0.25 km (S, day) / −0.95 km (N, night)", "Strong"),
        ("Water-cloud COT", "ACM-CAP", "median −4.8 (passive τ saturates ~5–8)", "Regime-limited"),
        ("Ice-cloud COT", "A-EBD", "median +2.0 (mean +7 skewed)", "Weak corr."),
        ("Effective radius", "ACM-CAP", "median +1 µm (robust bias, no skill)", "Robust bias"),
        ("Liquid water path", "ACM-CAP LWP", "median −16 g m⁻² (−34%)", "Inherits τ"),
        ("Cloud mask", "A-TC", "POD 0.69 · FAR 0.11 (misses = thin cirrus)", "Conservative"),
        ("Cloud phase", "A-TC", "POD_liq 90% · POD_ice 62% (liquid bias)", "Liquid-biased"),
    ]
    nr, nc = len(rows), 4
    tbl = s.shapes.add_table(nr, nc, Inches(0.55), Inches(1.55),
                             Inches(12.2), Inches(5.0)).table
    tbl.columns[0].width = Inches(2.7); tbl.columns[1].width = Inches(1.9)
    tbl.columns[2].width = Inches(5.4); tbl.columns[3].width = Inches(2.2)
    for ci in range(nc):
        for ri in range(nr):
            cell = tbl.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            run = para.add_run(); run.text = rows[ri][ci]
            run.font.size = Pt(14 if ri == 0 else 13)
            run.font.bold = (ri == 0 or ci == 0)
            if ri == 0:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = (RGBColor(0xEF, 0xF3, 0xF8) if ri % 2
                                            else RGBColor(0xFF, 0xFF, 0xFF))
    set_notes(s,
        "A one-look summary of every variable. The green rows are trustworthy - cloud-top "
        "height and effective radius, plus a conservative cloud mask. The amber rows are "
        "regime-limited - the optical depths and the water path, all tied to the same "
        "saturation. The red row is the phase, our clearest weakness. If you take one thing "
        "from this table: the thermal and microphysical-shape products are solid; the solar "
        "optical-depth magnitude and the phase are where the work is.")

    # ---- 10. Summary & next steps ----
    s = blank(prs)
    title_bar(s, "Summary — what is validated, and what is next",
              "Every EarthCARE-validatable dimension of the Dec-2025 sample is complete")
    bullets(s, [
        ("VALIDATED (all median-primary, with confidence intervals):", 0, True),
        ("CTH · water-COT · ice-COT · CER · CWP · cloud mask · phase — plus collocation methodology, surface-type & phase stratification.", 1, False),
        ("THE UNIFYING STORY — two independent limitations of the polar-daytime solar retrieval:", 0, True),
        ("(1) bright-surface optical-depth SATURATION (surface-driven; propagates into CWP; spares open water),", 1, False),
        ("(2) an intrinsic LIQUID PHASE BIAS (POD_ice 62%; surface-independent).", 1, False),
        ("Thermal CTH is strong (−0.25 km); CER is robust. These are the trustworthy products over the cryosphere.", 1, False),
        ("NEXT STEPS (require new processing, not further analysis):", 0, True),
        ("Arctic / boreal-summer month (Arctic sea-ice & Greenland) · Sentinel-3B · other seasons.", 1, False),
    ], left_in=0.7, top_in=1.5, width_in=12.0, height_in=5.6, size=17)
    set_notes(s,
        "To close. Everything is validated and defensible as an Antarctic-summer daytime "
        "study, with cloud-top height also bi-hemispheric. The unifying message is two "
        "independent limitations: a surface-driven optical-depth saturation that also hits "
        "the water path, and an intrinsic liquid phase bias. The reliable products are "
        "cloud-top height and effective radius. The next steps all require new processing "
        "rather than more analysis - a boreal-summer month to test the Arctic, Sentinel-3B "
        "for a second platform, and other seasons. Happy to take questions.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("wrote", OUT, f"({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
